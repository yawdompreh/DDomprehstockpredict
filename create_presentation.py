from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Create a PowerPoint presentation
presentation = Presentation()

# Title Slide
slide_title = presentation.slides.add_slide(presentation.slide_layouts[0])

# Set title and subtitle
title = slide_title.shapes.title
subtitle = slide_title.placeholders[1]

title.text = "AI-Powered Development Stack for Enterprise Banking"
subtitle.text = "A Comprehensive Overview"

# Function to add a slide with a title and content
def add_slide(title_text, content_text):
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])  # Title and Content layout
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = title_text
    content.text = content_text

# Add detailed slides (An example format with general placeholders)
add_slide("Technical Architecture", "Describe the technical architecture here.")
add_slide("Banking Implementation", "Outline banking implementation strategies here.")
add_slide("Knowledge Graphs", "Explain the importance of knowledge graphs in banking.")
add_slide("Practical Examples", "Include real-life examples of AI in banking.")
add_slide("Adoption Roadmap", "Provide a strategic roadmap for adoption.")
add_slide("Appendices", "Include appendices and supplementary information here.")

# Slide formatting - adjust font size, alignment, and styles as needed
for slide in presentation.slides:
    for shape in slide.shapes:
        if hasattr(shape, 'text_frame'):
            for paragraph in shape.text_frame.paragraphs:
                paragraph.font.size = Pt(18)  # Set font size
                paragraph.alignment = PP_ALIGN.LEFT  # Align text to the left

# Save the presentation
presentation.save('AI_Powered_Development_Stack_for_Enterprise_Banking.pptx')
